#!/usr/bin/env python3
"""Build .docx, .pptx, .pdf v1 from exports/trestle-brief-doc.html.

soffice was unavailable in this environment, so we build natively:
- weasyprint   -> PDF (renders the HTML as-is, full styling preserved)
- python-docx  -> Word
- python-pptx  -> PowerPoint (one section heading per slide)
"""
from pathlib import Path
import re
from bs4 import BeautifulSoup, NavigableString
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from weasyprint import HTML

HERE = Path(__file__).parent
SRC = HERE / "trestle-brief-doc.html"
DOCX_OUT = HERE / "trestle-brief.docx"
PPTX_OUT = HERE / "trestle-brief.pptx"
PDF_OUT = HERE / "trestle-brief.pdf"


def _inline_text(node):
    """Flatten an element's text content, collapsing whitespace."""
    if node is None:
        return ""
    text = node.get_text(" ", strip=True)
    return re.sub(r"\s+", " ", text)


# ---------- DOCX ----------

DOCX_BLUE = RGBColor(0x2C, 0x5A, 0x8A)
DOCX_AMBER = RGBColor(0xB0, 0x7B, 0x3A)
DOCX_GREY = RGBColor(0x6A, 0x73, 0x7C)


def _style_heading(p, size_pt, color, bold=True, upper=False):
    for run in p.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        if upper:
            run.text = run.text.upper()


def build_docx(soup: BeautifulSoup, out: Path):
    doc = Document()

    # Default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    body = soup.body
    for el in body.children:
        if isinstance(el, NavigableString):
            continue
        name = el.name
        if name == "h1":
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 22, RGBColor(0x1A, 0x1F, 0x25))
        elif name == "h2":
            doc.add_paragraph()  # spacer
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 14, DOCX_BLUE, upper=True)
        elif name == "h3":
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 12, DOCX_AMBER)
        elif name == "h4":
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 10, RGBColor(0x2A, 0x30, 0x38), upper=True)
        elif name == "p":
            cls = (el.get("class") or [])
            text = _inline_text(el)
            if not text:
                continue
            p = doc.add_paragraph(text)
            if "subtitle" in cls or "muted" in cls:
                for run in p.runs:
                    run.font.color.rgb = DOCX_GREY
                    run.font.size = Pt(10)
        elif name == "ul":
            for li in el.find_all("li", recursive=False):
                doc.add_paragraph(_inline_text(li), style="List Bullet")
        elif name == "table":
            rows = el.find_all("tr")
            if not rows:
                continue
            cols = max(len(r.find_all(["td", "th"])) for r in rows)
            table = doc.add_table(rows=0, cols=cols)
            table.style = "Light Grid Accent 1"
            for r in rows:
                cells = r.find_all(["th", "td"])
                row = table.add_row().cells
                for i, c in enumerate(cells):
                    if i < cols:
                        row[i].text = _inline_text(c)
        elif name == "div":
            cls = el.get("class") or []
            if "page-break" in cls:
                doc.add_page_break()
                continue
            if "card" in cls:
                # Render heading + child paragraphs/lists with a leading marker
                _render_card_docx(doc, el)
            elif "grid2" in cls:
                # Render each child card sequentially
                for child in el.find_all("div", class_="card", recursive=False):
                    _render_card_docx(doc, child)
            else:
                _render_card_docx(doc, el)

    doc.save(str(out))


def _render_card_docx(doc, card):
    """Render a card-like container: headings, paragraphs, kv grids, gauges, lists."""
    for el in card.children:
        if isinstance(el, NavigableString):
            continue
        name = el.name
        cls = el.get("class") or []
        if name == "h3":
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 12, DOCX_AMBER)
        elif name == "h4":
            p = doc.add_paragraph(_inline_text(el))
            _style_heading(p, 10, RGBColor(0x2A, 0x30, 0x38), upper=True)
        elif name == "p":
            text = _inline_text(el)
            if text:
                doc.add_paragraph(text)
        elif name == "ul":
            for li in el.find_all("li", recursive=False):
                doc.add_paragraph(_inline_text(li), style="List Bullet")
        elif name == "div" and "kv" in cls:
            _render_kv_docx(doc, el)
        elif name == "div" and "gauges" in cls:
            _render_gauges_docx(doc, el)
        elif name == "table":
            rows = el.find_all("tr")
            if rows:
                cols = max(len(r.find_all(["td", "th"])) for r in rows)
                table = doc.add_table(rows=0, cols=cols)
                table.style = "Light Grid Accent 1"
                for r in rows:
                    cells = r.find_all(["th", "td"])
                    row = table.add_row().cells
                    for i, c in enumerate(cells):
                        if i < cols:
                            row[i].text = _inline_text(c)


def _render_kv_docx(doc, kv):
    """Render a kv grid as a 2-column borderless table."""
    items = []
    children = [c for c in kv.children if not isinstance(c, NavigableString)]
    i = 0
    while i < len(children) - 1:
        k = children[i]
        v = children[i + 1]
        if "k" in (k.get("class") or []):
            items.append((_inline_text(k), _inline_text(v)))
            i += 2
        else:
            i += 1
    if not items:
        return
    table = doc.add_table(rows=0, cols=2)
    table.autofit = False
    for k, v in items:
        row = table.add_row().cells
        row[0].text = k.upper()
        row[1].text = v
        row[0].width = Inches(1.2)
        for p in row[0].paragraphs:
            for run in p.runs:
                run.font.size = Pt(8.5)
                run.font.color.rgb = DOCX_GREY
                run.font.bold = True


def _render_gauges_docx(doc, gauges):
    """Render gauge rows as plain text lines."""
    for row in gauges.find_all("div", class_="gauge-row"):
        label = row.find("span", class_="gauge-label")
        rest_parts = []
        for span in row.find_all("span"):
            if span is label:
                continue
            t = _inline_text(span)
            if t:
                rest_parts.append(t)
        if not label:
            continue
        line = f"{_inline_text(label)}: {' '.join(rest_parts)}"
        p = doc.add_paragraph(line)
        for run in p.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = DOCX_GREY


# ---------- PPTX ----------

PPTX_BLUE = PRGBColor(0x2C, 0x5A, 0x8A)
PPTX_DARK = PRGBColor(0x1A, 0x1F, 0x25)
PPTX_GREY = PRGBColor(0x6A, 0x73, 0x7C)


def _add_title_slide(prs, title, subtitle=""):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(PInches(0.5), PInches(1.5), PInches(9), PInches(1.2))
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.runs[0].font.size = PPt(34)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = PPTX_DARK
    if subtitle:
        sx = slide.shapes.add_textbox(PInches(0.5), PInches(2.7), PInches(9), PInches(2))
        stf = sx.text_frame
        stf.word_wrap = True
        stf.text = subtitle
        sr = stf.paragraphs[0].runs[0]
        sr.font.size = PPt(14)
        sr.font.color.rgb = PPTX_GREY


def _add_content_slide(prs, title, bullets):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title
    tx = slide.shapes.add_textbox(PInches(0.4), PInches(0.3), PInches(9.2), PInches(0.8))
    tf = tx.text_frame
    tf.text = title
    r = tf.paragraphs[0].runs[0]
    r.font.size = PPt(24)
    r.font.bold = True
    r.font.color.rgb = PPTX_BLUE
    # Body
    bx = slide.shapes.add_textbox(PInches(0.4), PInches(1.2), PInches(9.2), PInches(5.8))
    btf = bx.text_frame
    btf.word_wrap = True
    if not bullets:
        bullets = [""]
    for i, (level, text) in enumerate(bullets):
        if i == 0:
            para = btf.paragraphs[0]
        else:
            para = btf.add_paragraph()
        para.text = text
        para.level = min(level, 4)
        for run in para.runs:
            run.font.size = PPt(16 if level == 0 else 14)
            run.font.color.rgb = PPTX_DARK


def _flatten_section_to_bullets(section_nodes):
    """Convert a list of body nodes (between h2 boundaries) into pptx bullets."""
    bullets = []
    for el in section_nodes:
        if isinstance(el, NavigableString):
            continue
        name = el.name
        cls = el.get("class") or []
        if name == "h3":
            bullets.append((0, _inline_text(el).upper()))
        elif name == "h4":
            bullets.append((1, _inline_text(el)))
        elif name == "p":
            text = _inline_text(el)
            if text:
                bullets.append((1, text))
        elif name == "ul":
            for li in el.find_all("li", recursive=False):
                bullets.append((2, _inline_text(li)))
        elif name == "table":
            for r in el.find_all("tr"):
                cells = r.find_all(["th", "td"])
                bullets.append((1, " — ".join(_inline_text(c) for c in cells)))
        elif name == "div":
            if "card" in cls or "grid2" in cls or "kv" in cls or "gauges" in cls:
                bullets.extend(_card_to_bullets(el))
    return bullets


def _card_to_bullets(card):
    out = []
    for el in card.children:
        if isinstance(el, NavigableString):
            continue
        name = el.name
        cls = el.get("class") or []
        if name == "h3":
            out.append((0, _inline_text(el)))
        elif name == "h4":
            out.append((1, _inline_text(el)))
        elif name == "p":
            text = _inline_text(el)
            if text:
                out.append((1, text))
        elif name == "ul":
            for li in el.find_all("li", recursive=False):
                out.append((2, _inline_text(li)))
        elif name == "table":
            for r in el.find_all("tr"):
                cells = r.find_all(["th", "td"])
                out.append((1, " — ".join(_inline_text(c) for c in cells)))
        elif name == "div" and "kv" in cls:
            children = [c for c in el.children if not isinstance(c, NavigableString)]
            i = 0
            while i < len(children) - 1:
                k = children[i]
                v = children[i + 1]
                if "k" in (k.get("class") or []):
                    out.append((1, f"{_inline_text(k)}: {_inline_text(v)}"))
                    i += 2
                else:
                    i += 1
        elif name == "div" and "gauges" in cls:
            for row in el.find_all("div", class_="gauge-row"):
                label = row.find("span", class_="gauge-label")
                rest_parts = [
                    _inline_text(s)
                    for s in row.find_all("span")
                    if s is not label and _inline_text(s)
                ]
                if label:
                    out.append((2, f"{_inline_text(label)}: {' '.join(rest_parts)}"))
        elif name == "div" and ("card" in cls or "grid2" in cls):
            out.extend(_card_to_bullets(el))
    return out


def build_pptx(soup: BeautifulSoup, out: Path):
    prs = Presentation()
    prs.slide_width = PInches(13.33)
    prs.slide_height = PInches(7.5)

    body = soup.body
    h1 = body.find("h1")
    sub = body.find("p", class_="subtitle")
    _add_title_slide(prs, _inline_text(h1), _inline_text(sub) if sub else "")

    # Group nodes by h2 sections
    sections = []
    current_title = None
    current_nodes = []
    for el in body.children:
        if isinstance(el, NavigableString):
            continue
        if el.name == "h2":
            if current_title is not None:
                sections.append((current_title, current_nodes))
            current_title = _inline_text(el)
            current_nodes = []
        elif el.name in ("h1",) or (el.name == "p" and "subtitle" in (el.get("class") or [])):
            continue
        else:
            if current_title is not None:
                current_nodes.append(el)
    if current_title is not None:
        sections.append((current_title, current_nodes))

    # For each section: one slide for the section overview, plus a slide per inner h3 (cards / deep dives)
    for title, nodes in sections:
        # Split nodes by their own h3 boundaries to make per-card slides for big sections
        h3_groups = []
        cur_h3 = None
        cur_nodes = []
        # Find h3s at the top level OR nested inside cards
        for n in nodes:
            if n.name == "h3":
                if cur_h3 is not None:
                    h3_groups.append((cur_h3, cur_nodes))
                cur_h3 = _inline_text(n)
                cur_nodes = []
            elif n.name == "div" and "card" in (n.get("class") or []):
                # Card with its own h3 inside
                inner_h3 = n.find("h3", recursive=False)
                if inner_h3:
                    if cur_h3 is not None:
                        h3_groups.append((cur_h3, cur_nodes))
                    cur_h3 = _inline_text(inner_h3)
                    # Strip the heading so it isn't duplicated
                    cur_nodes = [n]
                else:
                    cur_nodes.append(n)
            elif n.name == "div" and "grid2" in (n.get("class") or []):
                # Multiple cards inline — emit each as its own group
                if cur_h3 is not None:
                    h3_groups.append((cur_h3, cur_nodes))
                    cur_h3 = None
                    cur_nodes = []
                for sub_card in n.find_all("div", class_="card", recursive=False):
                    inner_h3 = sub_card.find("h3", recursive=False)
                    h3_groups.append((_inline_text(inner_h3) if inner_h3 else title, [sub_card]))
            else:
                cur_nodes.append(n)
        if cur_h3 is not None or cur_nodes:
            h3_groups.append((cur_h3 or title, cur_nodes))

        # Section divider slide
        _add_content_slide(prs, title, [(0, "Section overview")] +
                           [(1, h or "") for h, _ in h3_groups if h])

        for sub_title, sub_nodes in h3_groups:
            bullets = _flatten_section_to_bullets(sub_nodes)
            if not bullets:
                continue
            # Chunk slides if too many bullets
            chunk = 14
            for i in range(0, len(bullets), chunk):
                chunk_title = sub_title
                if i > 0:
                    chunk_title = f"{sub_title} (cont.)"
                _add_content_slide(prs, chunk_title, bullets[i:i + chunk])

    prs.save(str(out))


# ---------- PDF (weasyprint) ----------

def build_pdf(src: Path, out: Path):
    HTML(filename=str(src)).write_pdf(str(out))


def main():
    soup = BeautifulSoup(SRC.read_text(), "html.parser")
    print("Building DOCX...")
    build_docx(soup, DOCX_OUT)
    print(f"  -> {DOCX_OUT}")
    print("Building PPTX...")
    build_pptx(soup, PPTX_OUT)
    print(f"  -> {PPTX_OUT}")
    print("Building PDF...")
    build_pdf(SRC, PDF_OUT)
    print(f"  -> {PDF_OUT}")


if __name__ == "__main__":
    main()
